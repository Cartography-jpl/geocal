from geocal_swig import *
try:
    import matplotlib.pyplot as plt
    from matplotlib.backends.backend_pdf import PdfPages
    from pptx import Presentation
    from pptx.util import Inches
    import numpy as np
    import cv2
except ImportError:
    pass
import pickle
from test_support import *
from io import BytesIO

# We can't actually run this, it depends on test data (from the Mars Helicopter
# testing) and a version of opencv with SIFT etc in it. However, this
# is a great example of producing matplotlib graphs as both PDF and a
# PowerPoint presentation. So we include this code here as an example,
# even though we can't actually run this as part of our normal tests
@pytest.fixture(scope="module")
def pdf_file():
    pdf = PdfPages("test_plots.pdf")
    yield pdf
    pdf.close()

@pytest.fixture(scope="module")
def ppt_file():
    ppt = Presentation()
    yield ppt
    ppt.save("test_plots.ppt")
    
@pytest.fixture()
def igccol():
    dem = geocal.SrtmDem()
    cam = geocal.SimpleCamera()
    tm = geocal.Time.parse_time("2014-10-02T12:00:00Z")
    orb_data = geocal.KeplerOrbit().orbit_data(tm)
    img1 = geocal.VicarLiteRasterImage("052.VIC")
    igc1 = geocal.OrbitDataImageGroundConnection(orb_data, cam, dem, img1,
                                                 "Image 052")
    img2 = geocal.VicarLiteRasterImage("056.VIC")
    igc2 = geocal.OrbitDataImageGroundConnection(orb_data, cam, dem, img2,
                                                 "Image 056")
    igccol = geocal.IgcArray([igc1, igc2])
    return igccol

def ppt_save(ppt_file):
    blank_slide_layout = ppt_file.slide_layouts[6]
    slide = ppt_file.slides.add_slide(blank_slide_layout)
    img = BytesIO()
    plt.savefig(img, dpi=300)
    pic = slide.shapes.add_picture(img, Inches(0), Inches(0), height=Inches(7.5))
    
def plot_data(igccol, tplist, title):
    plt.clf()
    plt.suptitle(title)
    for i in range(2):
        plt.subplot(1, 2, i+1)
        plt.title(igccol.title(i))
        d = igccol.image(i).read_all()
        max = np.max(d)
        if(max > 0):
            min = np.min(d[d > 0])
        else:
            min = 0
        for tp in tplist:
            ic = tp.image_coordinate(i)
            plt.plot(ic.sample, ic.line, 'r+')
        plt.imshow(d, cmap = plt.cm.gray, vmin = min, vmax = max)
        frame = plt.gca()
        frame.axes.get_xaxis().set_visible(False)
        frame.axes.get_yaxis().set_visible(False)
    plt.tight_layout()

@skip    
def test_sift_default(pdf_file, ppt_file, igccol):
    number_feature = 500
    number_octave_levels = 4
    fdetect = cv2.xfeatures2d.SIFT_create(number_feature, number_octave_levels)
    tp_gen = geocal.TiePointCollectFM(igccol, fdetect=fdetect,
                                      skip_ray_intersect=True)
    tplist = tp_gen.tie_point_list()
    print("SIFT - %d points" % len(tplist))
    plot_data(igccol, tplist, "SIFT default\n%d points" % len(tplist))
    ppt_save(ppt_file)
    pdf_file.savefig()

@skip    
def test_surf_default(pdf_file, ppt_file, igccol):
    hessian_threshold=400
    fdetect = cv2.xfeatures2d.SURF_create(hessian_threshold)
    tp_gen = geocal.TiePointCollectFM(igccol, fdetect=fdetect,
                                      skip_ray_intersect=True)
    tplist = tp_gen.tie_point_list()
    print("Surf - %d points" % len(tplist))
    plot_data(igccol, tplist, "SURF default\n%d points" % len(tplist))
    ppt_save(ppt_file)
    pdf_file.savefig()

@skip    
def test_kaze_default(pdf_file, ppt_file, igccol):
    fdetect = cv2.KAZE.create()
    tp_gen = geocal.TiePointCollectFM(igccol, fdetect=fdetect,
                                      skip_ray_intersect=True)
    tplist = tp_gen.tie_point_list()
    print("KAZE - %d points" % len(tplist))
    plot_data(igccol, tplist, "KAZE default\n%d points" % len(tplist))
    ppt_save(ppt_file)
    pdf_file.savefig()

@skip    
def test_akaze_default(pdf_file, ppt_file, igccol):
    fdetect = cv2.AKAZE.create()
    tp_gen = geocal.TiePointCollectFM(igccol, fdetect=fdetect,
                                      norm_type = cv2.NORM_HAMMING,
                                      skip_ray_intersect=True)
    tplist = tp_gen.tie_point_list()
    print("AKAZE - %d points" % len(tplist))
    plot_data(igccol, tplist, "AKAZE default\n%d points" % len(tplist))
    ppt_save(ppt_file)
    pdf_file.savefig()

@skip    
def test_brisk_default(pdf_file, ppt_file, igccol):
    thresh = 30
    fdetect = cv2.BRISK.create(thresh)
    tp_gen = geocal.TiePointCollectFM(igccol, fdetect=fdetect,
                                      norm_type = cv2.NORM_HAMMING,
                                      skip_ray_intersect=True)
    tplist = tp_gen.tie_point_list()
    print("BRISK - %d points" % len(tplist))
    plot_data(igccol, tplist, "BRISK default\n%d points" % len(tplist))
    ppt_save(ppt_file)
    pdf_file.savefig()

    
@skip    
def test_orb_default(pdf_file, ppt_file, igccol):
    number_feature = 500
    fdetect = cv2.ORB.create(number_feature)
    tp_gen = geocal.TiePointCollectFM(igccol, fdetect=fdetect,
                                      norm_type = cv2.NORM_HAMMING,
                                      skip_ray_intersect=True)
    tplist = tp_gen.tie_point_list()
    print("ORB - %d points" % len(tplist))
    plot_data(igccol, tplist, "ORB default\n%d points" % len(tplist))
    ppt_save(ppt_file)
    pdf_file.savefig()

    
